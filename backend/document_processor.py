import hashlib
import io
from pathlib import Path
from typing import List, Dict, Any, Tuple
import pandas as pd
import fitz  # PyMuPDF
import docx
from pptx import Presentation

def calculate_file_hash(file_bytes: bytes) -> str:
    """Compute SHA-256 hash of file content."""
    sha256 = hashlib.sha256()
    sha256.update(file_bytes)
    return sha256.hexdigest()

def chunk_text(text: str, max_chars: int = 1200, overlap: int = 200) -> List[str]:
    """Helper to split a long string into overlapping chunks."""
    if len(text) <= max_chars:
        return [text]
    
    chunks = []
    start = 0
    while start < len(text):
        end = start + max_chars
        # Try to find a space or newline to split cleanly
        if end < len(text):
            # Search backwards for a space or punctuation
            split_idx = -1
            for i in range(end, max_chars // 2, -1):
                if text[i] in ['\n', '.', '?', '!', ' ', ';']:
                    split_idx = i + 1
                    break
            if split_idx != -1:
                end = split_idx
        
        chunks.append(text[start:end].strip())
        start = end - overlap
        if start < 0 or start >= len(text):
            break
            
    return chunks

def parse_pdf(file_path: Path) -> List[Dict[str, Any]]:
    """Parse PDF file using PyMuPDF and extract text per page."""
    chunks = []
    doc = fitz.open(file_path)
    
    # Try to extract the outline/table of contents to find section headings
    toc = doc.get_toc() # List of [level, title, page]
    
    def get_section_for_page(p_num: int) -> str:
        current_heading = "General"
        # TOC pages are usually 1-indexed, loop to find closest heading
        for level, title, page in toc:
            if page <= p_num:
                current_heading = title
        return current_heading

    for page_idx in range(len(doc)):
        page = doc[page_idx]
        page_num_str = f"Page {page_idx + 1}"
        section_heading = get_section_for_page(page_idx + 1)
        text = page.get_text()
        
        if not text.strip():
            continue
            
        # Split page text if it exceeds maximum chunk size
        text_chunks = chunk_text(text, max_chars=1200, overlap=150)
        for i, chunk in enumerate(text_chunks):
            chunks.append({
                "text": chunk,
                "page_number": page_num_str,
                "section_heading": section_heading
            })
            
    return chunks

def parse_docx(file_path: Path) -> List[Dict[str, Any]]:
    """Parse DOCX file, detecting headings and chunking paragraphs."""
    chunks = []
    doc = docx.Document(file_path)
    
    current_section = "Introduction"
    current_text_block = []
    block_char_count = 0
    para_index = 0
    
    # We will accumulate paragraphs until they hit a threshold, then output a chunk.
    # If we hit a heading, we output the previous chunk and update the heading.
    for para in doc.paragraphs:
        text = para.text.strip()
        if not text:
            continue
            
        # Check if paragraph is a heading (Style names usually contain 'Heading' or 'Title')
        is_heading = para.style.name.startswith("Heading") or para.style.name == "Title"
        
        if is_heading:
            # Save accumulated text as a chunk
            if current_text_block:
                full_text = "\n".join(current_text_block)
                for c in chunk_text(full_text, max_chars=1200, overlap=150):
                    chunks.append({
                        "text": c,
                        "page_number": f"Paragraph {para_index - len(current_text_block) + 1}",
                        "section_heading": current_section
                    })
                current_text_block = []
                block_char_count = 0
                
            current_section = text
        else:
            current_text_block.append(text)
            block_char_count += len(text)
            para_index += 1
            
            if block_char_count >= 1000:
                full_text = "\n".join(current_text_block)
                for c in chunk_text(full_text, max_chars=1200, overlap=150):
                    chunks.append({
                        "text": c,
                        "page_number": f"Paragraph {para_index - len(current_text_block) + 1}",
                        "section_heading": current_section
                    })
                current_text_block = []
                block_char_count = 0
                
    # Add any remaining text
    if current_text_block:
        full_text = "\n".join(current_text_block)
        for c in chunk_text(full_text, max_chars=1200, overlap=150):
            chunks.append({
                "text": c,
                "page_number": f"Paragraph {para_index - len(current_text_block) + 1}",
                "section_heading": current_section
            })
            
    # Fallback if no chunks were extracted but file has tables
    for table_idx, table in enumerate(doc.tables):
        table_data = []
        for row in table.rows:
            row_data = [cell.text.strip() for cell in row.cells]
            table_data.append(" | ".join(row_data))
        if table_data:
            table_text = f"Table {table_idx + 1}:\n" + "\n".join(table_data)
            chunks.append({
                "text": table_text,
                "page_number": f"Table {table_idx + 1}",
                "section_heading": "Tables"
            })
            
    return chunks

def parse_pptx(file_path: Path) -> List[Dict[str, Any]]:
    """Parse PPTX file and extract text slide-by-slide."""
    chunks = []
    prs = Presentation(file_path)
    
    for slide_idx, slide in enumerate(prs.slides):
        slide_num_str = f"Slide {slide_idx + 1}"
        slide_text_elements = []
        slide_title = f"Slide {slide_idx + 1}"
        
        # Try to find a title shape
        if slide.shapes.title:
            slide_title = slide.shapes.title.text.strip()
            
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                # Avoid duplicate slide title if we already extracted it
                if shape == slide.shapes.title:
                    continue
                slide_text_elements.append(shape.text.strip())
                
        text_content = "\n".join(slide_text_elements)
        if not text_content.strip() and not slide.shapes.title:
            continue
            
        full_slide_text = f"Title: {slide_title}\n{text_content}"
        
        # Split slide text if it's exceptionally long
        text_chunks = chunk_text(full_slide_text, max_chars=1200, overlap=150)
        for i, chunk in enumerate(text_chunks):
            chunks.append({
                "text": chunk,
                "page_number": slide_num_str,
                "section_heading": slide_title
            })
            
    return chunks

def parse_xlsx(file_path: Path) -> List[Dict[str, Any]]:
    """Parse XLSX Excel files, processing each sheet."""
    chunks = []
    try:
        # Load all sheets
        excel_file = pd.ExcelFile(file_path)
        for sheet_name in excel_file.sheet_names:
            df = pd.read_excel(file_path, sheet_name=sheet_name)
            if df.empty:
                continue
                
            # Convert sheet into markdown-like text representation or CSV format
            # Process in groups of rows to make reasonable chunks
            rows_per_chunk = 30
            num_rows = len(df)
            
            for start_row in range(0, num_rows, rows_per_chunk):
                end_row = min(start_row + rows_per_chunk, num_rows)
                sub_df = df.iloc[start_row:end_row]
                
                # Convert this slice to string representation
                csv_buffer = io.StringIO()
                sub_df.to_csv(csv_buffer, index=False)
                sheet_text = f"Sheet: {sheet_name} (Rows {start_row + 1} - {end_row})\n" + csv_buffer.getvalue()
                
                chunks.append({
                    "text": sheet_text,
                    "page_number": f"Sheet {sheet_name}, Rows {start_row + 1}-{end_row}",
                    "section_heading": f"Sheet: {sheet_name}"
                })
    except Exception as e:
        # Fallback to simple pandas read
        df = pd.read_excel(file_path)
        csv_buffer = io.StringIO()
        df.to_csv(csv_buffer, index=False)
        chunks.append({
            "text": csv_buffer.getvalue(),
            "page_number": "Sheet 1",
            "section_heading": "Excel Data"
        })
        
    return chunks

def parse_csv(file_path: Path) -> List[Dict[str, Any]]:
    """Parse CSV files, chunking rows."""
    chunks = []
    df = pd.read_csv(file_path)
    if df.empty:
        return []
        
    rows_per_chunk = 40
    num_rows = len(df)
    
    for start_row in range(0, num_rows, rows_per_chunk):
        end_row = min(start_row + rows_per_chunk, num_rows)
        sub_df = df.iloc[start_row:end_row]
        
        csv_buffer = io.StringIO()
        sub_df.to_csv(csv_buffer, index=False)
        csv_text = f"CSV Data (Rows {start_row + 1} - {end_row})\n" + csv_buffer.getvalue()
        
        chunks.append({
            "text": csv_text,
            "page_number": f"Rows {start_row + 1}-{end_row}",
            "section_heading": "Data Rows"
        })
        
    return chunks

def parse_txt(file_path: Path) -> List[Dict[str, Any]]:
    """Parse plain text files, chunking by paragraph/lines."""
    chunks = []
    
    # Try different encodings
    encodings = ["utf-8", "latin-1", "cp1252"]
    text = ""
    for enc in encodings:
        try:
            with open(file_path, "r", encoding=enc) as f:
                text = f.read()
            break
        except UnicodeDecodeError:
            continue
            
    if not text:
        return []
        
    # Split text into overlapping chunks
    text_chunks = chunk_text(text, max_chars=1200, overlap=150)
    for i, chunk in enumerate(text_chunks):
        # Calculate line estimation
        line_estimate = f"Section {i + 1}"
        chunks.append({
            "text": chunk,
            "page_number": line_estimate,
            "section_heading": "Document Content"
        })
        
    return chunks

def process_document(file_path: Path, file_type: str) -> List[Dict[str, Any]]:
    """Dispatch file to corresponding parser based on file type."""
    file_type = file_type.lower().strip('.')
    if file_type == 'pdf':
        return parse_pdf(file_path)
    elif file_type in ['docx', 'doc']:
        return parse_docx(file_path)
    elif file_type in ['pptx', 'ppt']:
        return parse_pptx(file_path)
    elif file_type in ['xlsx', 'xls']:
        return parse_xlsx(file_path)
    elif file_type == 'csv':
        return parse_csv(file_path)
    elif file_type in ['txt', 'md']:
        return parse_txt(file_path)
    else:
        # Fallback to text parsing
        return parse_txt(file_path)
